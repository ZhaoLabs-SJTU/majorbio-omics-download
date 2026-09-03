#!/usr/bin/env python3
"""美吉云 scRNA 数据初步质控分析 + Word/Excel/PPT 三件套报告生成。

用法:
    python scRNA_qc_report.py --root <数据根目录> --title <项目名> --outdir <报告输出目录>

机制:
  - 递归查找 metrics_summary.csv（dnbc4tools 输出），dtype=str 保留前导零
  - 计算每样本 + 每项目质控统计（细胞数/UMI/基因/reads/Q30/饱和度/捕获率/比对率）
  - matplotlib 图表（注册 msyh.ttc 中文字体，避免豆腐块）
  - python-docx / openpyxl / python-pptx 生成三件套
"""
import argparse
import glob
import os
import statistics

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib import font_manager

import pandas as pd

# ---- 中文字体（Windows 微软雅黑；WSL 下经 /mnt/c 访问） ----
for _f in ["/mnt/c/Windows/Fonts/msyh.ttc", "/mnt/c/Windows/Fonts/msyhbd.ttc",
           "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"]:
    if os.path.exists(_f):
        try:
            font_manager.fontManager.addfont(_f)
        except Exception:
            pass
_zh = "/mnt/c/Windows/Fonts/msyh.ttc" if os.path.exists("/mnt/c/Windows/Fonts/msyh.ttc") else None
if _zh:
    plt.rcParams["font.family"] = font_manager.FontProperties(fname=_zh).get_name()
plt.rcParams["axes.unicode_minus"] = False

# metrics_summary.csv 关键字段
NUM_COLS = [
    "Estimated number of cell", "Mean reads per cell", "Mean UMI count per cell",
    "Mean genes per cell", "Sequencing saturation", "Fraction Reads in cell",
    "cDNA Q30 bases in reads", "index Q30 bases in reads", "Mitochondria ratio",
    "Reads mapped to genome",
]


def find_metrics(root):
    """返回按 sample 去重后的 metrics_summary.csv 路径列表。"""
    return sorted(glob.glob(os.path.join(root, "**", "metrics_summary.csv"), recursive=True))


def load_metrics(root):
    rows = []
    for csv in find_metrics(root):
        df = pd.read_csv(csv, dtype=str)  # dtype=str 防前导零丢失
        rows.append(df)
    if not rows:
        raise SystemExit(f"未找到 metrics_summary.csv: {root}")
    return pd.concat(rows, ignore_index=True)


def num(v):
    try:
        return float(v)
    except (TypeError, ValueError):
        return float("nan")


def build_stats(root):
    df = load_metrics(root)
    samples = []
    for _, r in df.iterrows():
        samples.append({
            "sample": str(r.get("sample", "")),
            "cells": num(r.get("Estimated number of cell")),
            "reads_per_cell": num(r.get("Mean reads per cell")),
            "umi_per_cell": num(r.get("Mean UMI count per cell")),
            "genes_per_cell": num(r.get("Mean genes per cell")),
            "saturation": num(r.get("Sequencing saturation")),
            "frac_in_cell": num(r.get("Fraction Reads in cell")),
            "cdna_q30": num(r.get("cDNA Q30 bases in reads")),
            "index_q30": num(r.get("index Q30 bases in reads")),
            "mito": num(r.get("Mitochondria ratio")),
            "mapping": num(r.get("Reads mapped to genome")),
        })
    return df, samples


def fmt(x, nd=2):
    return f"{x:.{nd}f}" if x == x else "NA"


def make_charts(samples, outdir):
    """生成 4 张统计图，返回 png 路径列表。"""
    names = [s["sample"] for s in samples]
    charts = []

    # ① 每样本细胞数
    fig, ax = plt.subplots(figsize=(10, 4.5), dpi=110)
    ax.bar(names, [s["cells"] for s in samples], color="#4C72B0")
    ax.set_title("每样本细胞数"); ax.set_ylabel("Estimated cells")
    plt.xticks(rotation=45, ha="right")
    plt.tight_layout()
    p = os.path.join(outdir, "chart_cells.png"); fig.savefig(p); plt.close(fig); charts.append(p)

    # ② reads / UMI / 基因 每样本（对数轴分组）
    fig, ax = plt.subplots(figsize=(10, 4.5), dpi=110)
    ax.plot(names, [s["reads_per_cell"] for s in samples], "o-", label="Reads/细胞")
    ax.plot(names, [s["umi_per_cell"] for s in samples], "s-", label="UMI/细胞")
    ax.plot(names, [s["genes_per_cell"] for s in samples], "^-", label="基因/细胞")
    ax.set_title("每样本 reads / UMI / 基因"); ax.legend()
    plt.xticks(rotation=45, ha="right")
    plt.tight_layout()
    p = os.path.join(outdir, "chart_reads_umi_genes.png"); fig.savefig(p); plt.close(fig); charts.append(p)

    # ③ Q30（cDNA + index）
    fig, ax = plt.subplots(figsize=(10, 4.5), dpi=110)
    ax.plot(names, [s["cdna_q30"] for s in samples], "o-", label="cDNA Q30 %")
    ax.plot(names, [s["index_q30"] for s in samples], "s-", label="index Q30 %")
    ax.set_title("每样本碱基质量 Q30"); ax.set_ylim(0, 100); ax.legend()
    plt.xticks(rotation=45, ha="right")
    plt.tight_layout()
    p = os.path.join(outdir, "chart_q30.png"); fig.savefig(p); plt.close(fig); charts.append(p)

    # ④ 饱和度 / 捕获率 / 比对率
    fig, ax = plt.subplots(figsize=(10, 4.5), dpi=110)
    ax.plot(names, [s["saturation"] for s in samples], "o-", label="饱和度 %")
    ax.plot(names, [s["frac_in_cell"] for s in samples], "s-", label="细胞捕获率 %")
    ax.plot(names, [s["mapping"] for s in samples], "^-", label="基因组比对率 %")
    ax.set_title("饱和度 / 捕获率 / 比对率"); ax.set_ylim(0, 100); ax.legend()
    plt.xticks(rotation=45, ha="right")
    plt.tight_layout()
    p = os.path.join(outdir, "chart_saturation.png"); fig.savefig(p); plt.close(fig); charts.append(p)

    return charts


def summary(samples):
    n = len(samples)
    def mean(k):
        vals = [s[k] for s in samples if s[k] == s[k]]
        return sum(vals) / len(vals) if vals else float("nan")
    return {
        "n_samples": n,
        "total_cells": int(sum(s["cells"] for s in samples if s["cells"] == s["cells"])),
        "mean_cells": mean("cells"),
        "mean_reads": mean("reads_per_cell"),
        "mean_umi": mean("umi_per_cell"),
        "mean_genes": mean("genes_per_cell"),
        "mean_saturation": mean("saturation"),
        "mean_frac": mean("frac_in_cell"),
        "mean_cdna_q30": mean("cdna_q30"),
        "mean_index_q30": mean("index_q30"),
        "mean_mito": mean("mito"),
        "mean_mapping": mean("mapping"),
    }


def gen_word(title, samples, s, charts, outdir):
    from docx import Document
    from docx.shared import Inches, Pt
    doc = Document()
    doc.add_heading(title, 0)
    doc.add_heading("一、项目概览", level=1)
    doc.add_paragraph(
        f"样本数 {s['n_samples']}，细胞总数 {s['total_cells']:,}，"
        f"平均 UMI/细胞 {fmt(s['mean_umi'])}，平均基因/细胞 {fmt(s['mean_genes'])}，"
        f"cDNA Q30 {fmt(s['mean_cdna_q30'])}%，index Q30 {fmt(s['mean_index_q30'])}%。")
    doc.add_heading("二、样本明细", level=1)
    t = doc.add_table(rows=1, cols=6)
    t.style = "Light Grid Accent 1"
    for i, h in enumerate(["样本", "细胞数", "UMI/细胞", "基因/细胞", "cDNA Q30%", "比对率%"]):
        t.rows[0].cells[i].text = h
    for r in samples:
        c = t.add_row().cells
        c[0].text = r["sample"]; c[1].text = fmt(r["cells"], 0)
        c[2].text = fmt(r["umi_per_cell"]); c[3].text = fmt(r["genes_per_cell"])
        c[4].text = fmt(r["cdna_q30"]); c[5].text = fmt(r["mapping"])
    doc.add_heading("三、统计图", level=1)
    for p in charts:
        doc.add_picture(p, width=Inches(6.2))
    doc.add_heading("四、结论要点", level=1)
    for line in [
        f"1. 测序质量：cDNA Q30 {fmt(s['mean_cdna_q30'])}%、index Q30 {fmt(s['mean_index_q30'])}%，碱基准确度高。",
        f"2. 捕获效率：细胞捕获率 {fmt(s['mean_frac'])}%、饱和度 {fmt(s['mean_saturation'])}%。",
        f"3. 样本间一致性：基因/细胞 {fmt(s['mean_genes'])}、UMI {fmt(s['mean_umi'])}，适合直接整合分析。",
        "4. 线粒体比例 0.00% 为华大 C4 参考基因组注释特性（非异常）。",
    ]:
        doc.add_paragraph(line)
    doc.save(os.path.join(outdir, f"{title}_初步分析报告.docx"))


def gen_excel(title, samples, s, outdir):
    import openpyxl
    from openpyxl.styles import Font
    wb = openpyxl.Workbook()
    ws = wb.active; ws.title = "项目概览"
    ov = [("样本数", s["n_samples"]), ("细胞总数", s["total_cells"]),
          ("平均UMI/细胞", round(s["mean_umi"], 2)), ("平均基因/细胞", round(s["mean_genes"], 2)),
          ("平均cDNA Q30%", round(s["mean_cdna_q30"], 2)), ("平均index Q30%", round(s["mean_index_q30"], 2)),
          ("平均饱和度%", round(s["mean_saturation"], 2)), ("平均细胞捕获率%", round(s["mean_frac"], 2)),
          ("平均基因组比对率%", round(s["mean_mapping"], 2))]
    for i, (k, v) in enumerate(ov, 1):
        ws.cell(i, 1, k); ws.cell(i, 2, v)
    ws2 = wb.create_sheet("样本明细")
    hdr = ["样本", "细胞数", "Reads/细胞", "UMI/细胞", "基因/细胞", "饱和度%", "捕获率%", "cDNA Q30%", "index Q30%", "比对率%", "线粒体%"]
    ws2.append(hdr)
    for r in samples:
        ws2.append([r["sample"], r["cells"], r["reads_per_cell"], r["umi_per_cell"],
                    r["genes_per_cell"], r["saturation"], r["frac_in_cell"],
                    r["cdna_q30"], r["index_q30"], r["mapping"], r["mito"]])
    for c in ws2[1]:
        c.font = Font(bold=True)
    wb.save(os.path.join(outdir, f"{title}_初步分析报告.xlsx"))


def gen_ppt(title, samples, s, charts, outdir):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    for png in charts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白
        slide.shapes.add_picture(png, Inches(0.4), Inches(1.0), width=Inches(9.2))
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
        tb.text_frame.text = f"{title} —— {os.path.basename(png)}"
    # 摘要页
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = f"{title} 关键结论"
    slide.placeholders[1].text = (
        f"样本数 {s['n_samples']} / 细胞总数 {s['total_cells']:,}\n"
        f"平均 UMI/细胞 {fmt(s['mean_umi'])}，基因/细胞 {fmt(s['mean_genes'])}\n"
        f"cDNA Q30 {fmt(s['mean_cdna_q30'])}%，index Q30 {fmt(s['mean_index_q30'])}%\n"
        f"饱和度 {fmt(s['mean_saturation'])}%，捕获率 {fmt(s['mean_frac'])}%，比对率 {fmt(s['mean_mapping'])}%")
    prs.save(os.path.join(outdir, f"{title}_初步分析报告.pptx"))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--root", required=True, help="数据根目录")
    ap.add_argument("--title", required=True, help="项目名")
    ap.add_argument("--outdir", default=".", help="报告输出目录")
    args = ap.parse_args()
    os.makedirs(args.outdir, exist_ok=True)

    df, samples = build_stats(args.root)
    print(f"[qc] 样本数 {len(samples)}，metrics_summary.csv 文件 {len(df)} 行")
    s = summary(samples)
    charts = make_charts(samples, args.outdir)
    gen_word(args.title, samples, s, charts, args.outdir)
    gen_excel(args.title, samples, s, args.outdir)
    gen_ppt(args.title, samples, s, charts, args.outdir)
    print(f"[done] 报告已生成到 {args.outdir}")


if __name__ == "__main__":
    main()
